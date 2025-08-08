from flask import Flask, request, jsonify
import pandas as pd
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
import tempfile
import os
import base64

app = Flask(__name__)

BACKGROUND_COLOR = RGBColor(230, 230, 230)
TITLE_FONT = 'Calibri'
BODY_FONT = 'Calibri'
TITLE_FONT_SIZE = Pt(28)
BODY_FONT_SIZE = Pt(18)
TEXT_COLOR = RGBColor(50, 50, 50)

def set_slide_background_gray(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BACKGROUND_COLOR

def add_logo(slide, logo_path, prs):
    left = prs.slide_width - Inches(1.5)
    top = Inches(0.2)
    height = Inches(1)
    slide.shapes.add_picture(logo_path, left, top, height=height)

def add_title_slide(prs, title, subtitle, logo_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background_gray(slide)
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), prs.slide_width - Inches(2), Inches(1))
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    p.font.name = TITLE_FONT
    p.font.size = TITLE_FONT_SIZE
    p.font.bold = True
    p.font.color.rgb = TEXT_COLOR

    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3), prs.slide_width - Inches(2), Inches(0.7))
    p2 = subtitle_box.text_frame.paragraphs[0]
    run2 = p2.add_run()
    run2.text = subtitle
    p2.font.name = BODY_FONT
    p2.font.size = Pt(16)
    p2.font.italic = True
    p2.font.color.rgb = TEXT_COLOR

    add_logo(slide, logo_path, prs)

def add_category_summary_slide(prs, category, df, logo_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background_gray(slide)
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), prs.slide_width - Inches(1.4), Inches(1))
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = category.upper()
    p.font.name = TITLE_FONT
    p.font.size = TITLE_FONT_SIZE
    p.font.bold = True
    p.font.color.rgb = TEXT_COLOR

    total_news = len(df)
    total_circulation = df['Circulação'].sum() if 'Circulação' in df.columns else 0

    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.7), prs.slide_width - Inches(1.4), Inches(1))
    p2 = content_box.text_frame.paragraphs[0]
    run2 = p2.add_run()
    run2.text = f"Total de Notícias: {total_news}\nTotal de Circulação: {total_circulation:,}".replace(",", ".")
    p2.font.name = BODY_FONT
    p2.font.size = BODY_FONT_SIZE
    p2.font.color.rgb = TEXT_COLOR

    add_logo(slide, logo_path, prs)

def add_news_slide(prs, title_text, circulation, logo_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background_gray(slide)
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.3), prs.slide_width - Inches(1.4), Inches(0.8))
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "NOTÍCIA"
    p.font.name = TITLE_FONT
    p.font.size = TITLE_FONT_SIZE
    p.font.bold = True
    p.font.color.rgb = TEXT_COLOR

    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), prs.slide_width - Inches(1.4), Inches(3))
    p2 = content_box.text_frame.paragraphs[0]
    run2 = p2.add_run()
    run2.text = title_text
    p2.font.name = BODY_FONT
    p2.font.size = BODY_FONT_SIZE
    p2.font.color.rgb = TEXT_COLOR

    circ_box = slide.shapes.add_textbox(Inches(0.7), Inches(4.6), prs.slide_width - Inches(1.4), Inches(0.7))
    p3 = circ_box.text_frame.paragraphs[0]
    run3 = p3.add_run()
    run3.text = f"Circulação: {circulation:,}".replace(",", ".")
    p3.font.name = BODY_FONT
    p3.font.size = BODY_FONT_SIZE
    p3.font.italic = True
    p3.font.color.rgb = TEXT_COLOR

    add_logo(slide, logo_path, prs)

@app.route('/gerar-pptx', methods=['POST'])
def gerar_pptx():
    if 'excel' not in request.files or 'logo' not in request.files:
        return jsonify({"erro": "Faltam os arquivos 'excel' e/ou 'logo'."}), 400

    excel_file = request.files['excel']
    logo_file = request.files['logo']

    with tempfile.TemporaryDirectory() as tmpdirname:
        excel_path = os.path.join(tmpdirname, "input.xlsx")
        logo_path = os.path.join(tmpdirname, "logo.png")
        pptx_path = os.path.join(tmpdirname, "output.pptx")

        excel_file.save(excel_path)
        logo_file.save(logo_path)

        try:
            xls = pd.ExcelFile(excel_path)
        except Exception as e:
            return jsonify({"erro": f"Erro ao ler Excel: {e}"}), 400

        sheets = [s for s in xls.sheet_names if not s.startswith("Sheet")]
        dataframes = {sheet: xls.parse(sheet) for sheet in sheets}

        prs = Presentation()
        add_title_slide(prs, "Relatório de Notícias", "Gerado automaticamente via API", logo_path)

        for category, df in dataframes.items():
            if df.empty or 'Título' not in df.columns:
                continue
            df = df.dropna(subset=['Título'])
            if df.empty:
                continue
            add_category_summary_slide(prs, category, df, logo_path)
            for _, row in df.iterrows():
                title = str(row['Título'])
                circulation = row.get('Circulação', 0)
                add_news_slide(prs, title, circulation, logo_path)

        prs.save(pptx_path)

        # Converter para Base64
        with open(pptx_path, "rb") as f:
            pptx_bytes = f.read()
            pptx_base64 = base64.b64encode(pptx_bytes).decode('utf-8')

        return jsonify({
            "filename": "relatorio.pptx",
            "mimetype": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "file_base64": pptx_base64
        })

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=10000)
